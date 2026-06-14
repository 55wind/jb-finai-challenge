# -*- coding: utf-8 -*-
"""시스템 아키텍처 단독 PPT (편집 가능 도형).

- 5계층 데이터 흐름 (사용자 -> 프론트엔드 -> 백엔드/AI Agent -> DB/외부연동)
- 모델은 무료 로고 이미지(Ollama, HuggingFace/BGE-m3)로 표기하고,
  각 모델을 실제로 쓰는 파이프라인 단계와 연결선으로 이어 표현.
- 준법 오토파일럿은 아이콘 + 순서도(루프 포함)로 표현.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import os

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = os.path.join(HERE, '..', 'docs', 'assets', 'arch')
ICO_OLLAMA = os.path.join(ASSET, 'ollama.png')
ICO_BGE = os.path.join(ASSET, 'bge_hf.png')
ICO_AUTO = os.path.join(ASSET, 'autopilot.png')

C = lambda h: RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
JB_BLUE = C('00529b'); JB_DARK = C('003a6e'); INK = C('1a1d21'); MUTED = C('5b6470')
WHITE = C('ffffff'); LINE = C('cdd5de'); TEALc = C('0f9d58'); PURPc = C('5b3bbf')
SLATE = C('5b6470'); AMBER = C('c78a06'); GREYc = C('6b7480')
GREY_T = C('eef1f5'); TEAL_T = C('e7f6ee'); PURP_T = C('efeaff')

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
s = prs.slides.add_slide(prs.slide_layouts[6]); I = Inches


def settxt(tf, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, m=0.05):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(m); tf.margin_right = Inches(m)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    f = True
    for ln in runs:
        p = tf.paragraphs[0] if f else tf.add_paragraph(); f = False
        p.alignment = align; p.space_after = Pt(0); p.space_before = Pt(0); p.line_spacing = 1.0
        for t, sz, b, col in ln:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = col; r.font.name = 'Malgun Gothic'


def shp(kind, x, y, w, h, runs, fill, line=LINE, lw=0.75, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, m=0.06):
    sp = s.shapes.add_shape(kind, x, y, w, h); sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    if runs: settxt(sp.text_frame, runs, align, anchor, m=m)
    return sp


def box(x, y, w, h, runs, fill, **k):
    return shp(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, runs, fill, **k)


def tbox(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); settxt(tb.text_frame, runs, align, anchor, m=0); return tb


def rarrow(x, y, w, h, lab=''):
    shp(MSO_SHAPE.RIGHT_ARROW, x, y, w, h, [], C('8aa0b5'), line=None)
    if lab:
        tbox(x - I(0.35), y - I(0.32), w + I(0.7), I(0.26), [[(lab, 9.5, True, MUTED)]], align=PP_ALIGN.CENTER)


def darrow(xc, y, h, lab=''):
    shp(MSO_SHAPE.DOWN_ARROW, xc - I(0.14), y, I(0.28), h, [], C('8aa0b5'), line=None)
    if lab:
        tbox(xc + I(0.2), y, I(2.5), h, [[(lab, 9.5, True, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)


def connect(x1, y1, x2, y2, color, w=1.75, arrow=False):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    if arrow:
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def pic(path, x, y, h):
    return s.shapes.add_picture(path, x, y, height=h)


# ===== 제목 =====
tbox(I(0.45), I(0.18), I(11), I(0.5), [[('JB 준법 코파일럿 - 시스템 아키텍처', 25, True, JB_DARK)]])
tbox(I(0.47), I(0.74), I(12.4), I(0.3),
     [[('사용자 · 프론트엔드 · 백엔드(AI Agent/모델) · DB/외부연동 - 전 과정 온프레미스 · 휴먼인더루프', 12, False, MUTED)]])
shp(MSO_SHAPE.RECTANGLE, I(0.45), I(1.10), I(12.43), Pt(2.5), [], JB_BLUE, line=None)

# ===== 왼쪽: 사용자 + 프론트엔드 (세로 스택) =====
lx, lw = I(0.45), I(2.5)
# 사용자
box(lx, I(1.5), lw, I(1.5), [], C('f4f6f8'), line=SLATE, lw=1.25)
tbox(lx, I(1.6), lw, I(0.3), [[('[ 사용자 ]', 13.5, True, SLATE)]], align=PP_ALIGN.CENTER)
box(lx + I(0.16), I(1.98), lw - I(0.32), I(0.45),
    [[('마케팅·지점 담당자', 10.5, True, INK), ('  콘텐츠 작성', 9, False, MUTED)]], WHITE, line=LINE, align=PP_ALIGN.LEFT, m=0.1)
box(lx + I(0.16), I(2.48), lw - I(0.32), I(0.45),
    [[('준법관리자', 10.5, True, INK), ('  검토·승인(HITL)', 9, False, MUTED)]], WHITE, line=LINE, align=PP_ALIGN.LEFT, m=0.1)
darrow(I(1.7), I(3.05), I(0.32))
# 프론트엔드
box(lx, I(3.45), lw, I(2.25), [], C('eafaf0'), line=TEALc, lw=1.25)
tbox(lx, I(3.55), lw, I(0.3), [[('[ 프론트엔드 / UI ]', 12.5, True, TEALc)]], align=PP_ALIGN.CENTER)
tbox(lx, I(3.87), lw, I(0.24), [[('정적 SPA · 바닐라 JS (FastAPI 서빙)', 8.5, False, MUTED)]], align=PP_ALIGN.CENTER)
for i, u in enumerate(['심의 워크스페이스', '오인 시뮬레이터 *', '준법관리자 승인 뷰', '감사로그 뷰']):
    box(lx + I(0.16), I(4.18) + i * I(0.37), lw - I(0.32), I(0.32),
        [[(u, 10, '*' in u, INK)]], WHITE, line=C('bfe6cf'), align=PP_ALIGN.CENTER, m=0.05)
rarrow(I(3.02), I(4.35), I(0.42), I(0.4), 'REST(JSON)')

# ===== 백엔드 (AI Agent) 코어 =====
cx, cy, cw, ch = I(3.6), I(1.5), I(9.28), I(4.3)
box(cx, cy, cw, ch, [], C('faf8ff'), line=PURPc, lw=1.25)
ix = cx + I(0.2); iw = cw - I(0.4)
tbox(ix, cy + I(0.12), iw, I(0.28), [[('[ 백엔드 / Backend ]  FastAPI (app/main.py) · 온프레미스', 13, True, JB_BLUE)]])
tbox(ix, cy + I(0.46), iw, I(0.26), [[('AI Agent · Orchestrator - 명시적 파이썬 파이프라인', 12, True, PURPc)]])

# 파이프라인 7단계
stages = [('(1)분류', '결정적', GREYc, GREY_T), ('(2)RAG', 'BGE-m3', TEALc, TEAL_T),
          ('(3)룰엔진', '결정적', GREYc, GREY_T), ('(4)LLM심의', 'Ollama', PURPc, PURP_T),
          ('(5)시뮬레이터', 'Ollama', PURPc, PURP_T), ('(6)리포트', '결정적', GREYc, GREY_T),
          ('(7,8)승인·감사', '결정적', GREYc, GREY_T)]
gap = I(0.05); scw = (int(iw) - int(gap) * (len(stages) - 1)) // len(stages)
sy = cy + I(0.78); shh = I(0.82)
centers = []
for i, (nm, tag, col, tfill) in enumerate(stages):
    x = ix + i * (scw + int(gap))
    box(x, sy, scw, shh, [[(nm, 10, True, INK)], [(' ', 4, False, INK)], [(tag, 9, True, col)]],
        tfill, line=col, align=PP_ALIGN.CENTER, m=0.03)
    centers.append(x + scw // 2)
    if i < len(stages) - 1:
        tbox(x + scw - I(0.02), sy, I(0.1), shh, [[('>', 11, True, MUTED)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
stage_bottom = sy + shh

# 모델 박스 (로고) + 연결선
mby = cy + I(1.98); mbh = I(0.66)
# BGE-m3 (HuggingFace) : (2)RAG 와 연결
bge_cx = centers[1]; bge_x = bge_cx - I(1.15); bge_w = I(2.3)
box(bge_x, mby, bge_w, mbh, [], TEAL_T, line=TEALc, lw=1.25)
pic(ICO_BGE, bge_x + I(0.12), mby + I(0.13), I(0.4))
tbox(bge_x + I(0.6), mby, bge_w - I(0.68), mbh,
     [[('BGE-m3', 10.5, True, TEALc)], [('임베딩 (HuggingFace)', 8.5, False, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
connect(bge_cx, mby, centers[1], stage_bottom, TEALc, 2.0)
# Ollama : (4)LLM심의·(5)시뮬레이터 와 연결
oll_cx = (centers[3] + centers[4]) // 2; oll_x = oll_cx - I(1.4); oll_w = I(2.8)
box(oll_x, mby, oll_w, mbh, [], PURP_T, line=PURPc, lw=1.25)
pic(ICO_OLLAMA, oll_x + I(0.12), mby + I(0.11), I(0.44))
tbox(oll_x + I(0.64), mby, oll_w - I(0.72), mbh,
     [[('Ollama · Qwen2.5-7B', 10.5, True, PURPc)], [('LLM 추론 (온프레미스)', 8.5, False, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
connect(oll_cx - I(0.3), mby, centers[3], stage_bottom, PURPc, 2.0)
connect(oll_cx + I(0.3), mby, centers[4], stage_bottom, PURPc, 2.0)

# ===== 준법 오토파일럿 (아이콘 + 순서도) =====
ay = cy + I(2.84)
pic(ICO_AUTO, ix, ay - I(0.02), I(0.42))
tbox(ix + I(0.5), ay, iw - I(0.5), I(0.28),
     [[('준법 오토파일럿', 12, True, PURPc), (' - 위반 초안을 통과까지 자율 개선 (Ollama 기반)', 9.5, False, MUTED)]])
# 순서도 (파이프라인처럼 박스 + 화살표 + 루프)
fy = cy + I(3.24); fhh = I(0.6)
flow = ['① 심의·점수', '② LLM 고쳐쓰기', '③ 재심의', '④ 통과 / 미수렴->사람']
fcol = [GREYc, PURPc, GREYc, JB_DARK]
fgap = I(0.46); fcw = (int(iw) - int(fgap) * (len(flow) - 1)) // len(flow)
fcenters = []
for i, fl in enumerate(flow):
    x = ix + i * (fcw + int(fgap))
    tint = PURP_T if fcol[i] == PURPc else (GREY_T if fcol[i] == GREYc else C('e9eef5'))
    box(x, fy, fcw, fhh, [[(fl, 10, True, fcol[i] if fcol[i] != GREYc else INK)]], tint, line=fcol[i], align=PP_ALIGN.CENTER, m=0.04)
    fcenters.append((x + fcw // 2, x, x + fcw))
    if i < len(flow) - 1:
        rarrow(x + fcw + I(0.06), fy + I(0.16), I(0.30), I(0.28))
# 루프백 화살표: ③ 재심의 -> ② 고쳐쓰기 (미통과 시 반복)
loop_y = fy + fhh + I(0.14)
c3 = fcenters[2][0]; c2 = fcenters[1][0]
connect(c3, fy + fhh, c3, loop_y, PURPc, 1.75)
connect(c3, loop_y, c2, loop_y, PURPc, 1.75)
connect(c2, loop_y, c2, fy + fhh, PURPc, 1.75, arrow=True)
tbox(c2 - I(0.5), loop_y - I(0.04), (c3 - c2) + I(1.0), I(0.24), [[('미통과 시 (2)~(3) 반복 (통과까지)', 9, True, PURPc)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===== 백엔드 -> DB / 외부연동 =====
darrow(I(6.0), I(5.85), I(0.3), '조회·기록')
darrow(I(11.4), I(5.85), I(0.3), '연동')
shp(MSO_SHAPE.CAN, I(4.6), I(6.18), I(1.8), I(1.12),
    [[('SQLite', 10.5, True, WHITE)], [('승인·감사로그', 8.5, False, C('dfe6f0'))], [('내규·배포', 8.5, False, C('dfe6f0'))]],
    JB_DARK, line=None)
shp(MSO_SHAPE.CAN, I(6.6), I(6.18), I(1.8), I(1.12),
    [[('규제 KB', 10.5, True, WHITE)], [('regulations.json', 8, False, C('dfe6f0'))], [('rules.yaml', 8, False, C('dfe6f0'))]],
    C('2a6099'), line=None)
shp(MSO_SHAPE.CLOUD, I(10.1), I(6.05), I(2.78), I(1.35), [], C('fff7e6'), line=AMBER, lw=1.25)
tbox(I(10.1), I(6.18), I(2.78), I(0.24), [[('외부 연동 (어댑터)', 10, True, C('8a6200'))]], align=PP_ALIGN.CENTER)
for i, (nm, fid) in enumerate([('법제처 OpenAPI', 'F14'), ('FCM·SNS·알림톡', 'F15'), ('코어뱅킹 API', 'F16')]):
    tbox(I(10.2), I(6.44) + i * I(0.27), I(2.6), I(0.24),
         [[(nm + '  ', 9, True, C('6b4e00')), (fid, 8.5, True, AMBER)]], align=PP_ALIGN.CENTER)

out = 'docs/submission/JB준법코파일럿_시스템아키텍처.pptx'
prs.save(out)
print('saved', os.path.getsize(out), 'bytes')
